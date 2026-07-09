from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io

ACCENT_BLUE = RGBColor(99, 102, 241)
ACCENT_CYAN = RGBColor(34, 211, 238)
ACCENT_EMERALD = RGBColor(52, 211, 153)
WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(15, 15, 30)
LIGHT_TEXT = RGBColor(220, 220, 240)
MUTED_TEXT = RGBColor(150, 150, 180)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=LIGHT_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, title, bullets, accent=ACCENT_BLUE):
    # Accent bar on left
    add_shape(slide, Inches(0.5), Inches(1.2), Inches(0.06), Inches(0.5), accent)
    # Title
    add_text_box(slide, Inches(0.8), Inches(1.1), Inches(8.5), Inches(0.7), title, size=28, bold=True, color=WHITE)
    # Bullets
    y = Inches(2.0)
    for b in bullets:
        dot = add_text_box(slide, Inches(0.8), y, Inches(0.2), Inches(0.4), "●", size=12, color=accent)
        add_text_box(slide, Inches(1.1), y, Inches(7.5), Inches(0.4), b, size=16, color=LIGHT_TEXT)
        y += Inches(0.55)

def build_pptx(slides_data: list) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for sd in slides_data:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        add_bg(slide)

        t = sd.get("type", "content")

        if t == "title":
            # Dark gradient overlay bar
            add_shape(slide, Inches(0), Inches(1.8), Inches(10), Inches(2.2), RGBColor(30, 30, 60))
            # Accent line
            add_shape(slide, Inches(4.2), Inches(1.6), Inches(1.6), Inches(0.04), ACCENT_CYAN)
            # Title
            add_text_box(slide, Inches(1), Inches(2.0), Inches(8), Inches(1.2), sd.get("heading", ""), size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            # Subtitle
            sub = sd.get("subheading", "")
            if sub:
                add_text_box(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.5), sub, size=16, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)
            # Footer
            add_text_box(slide, Inches(1), Inches(4.5), Inches(8), Inches(0.3), "VEDA | AI-Powered Research", size=10, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

        elif t == "content":
            add_bullet_slide(slide, sd.get("heading", ""), sd.get("bullets", []),
                           ACCENT_CYAN if sd.get("accent") == "cyan" else ACCENT_EMERALD if sd.get("accent") == "emerald" else ACCENT_BLUE)

        elif t == "section":
            add_shape(slide, Inches(0), Inches(2.0), Inches(10), Inches(1.6), RGBColor(20, 20, 45))
            add_shape(slide, Inches(0), Inches(2.0), Inches(10), Inches(0.04), accent=sd.get("accent", ACCENT_BLUE))
            add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(0.8), sd.get("heading", ""), size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
            sub = sd.get("subheading", "")
            if sub:
                add_text_box(slide, Inches(1), Inches(3.0), Inches(8), Inches(0.4), sub, size=14, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

        elif t == "quote":
            add_shape(slide, Inches(1.5), Inches(1.5), Inches(7), Inches(2.5), RGBColor(20, 20, 45))
            add_text_box(slide, Inches(2), Inches(1.7), Inches(6), Inches(1.5), f"\"{sd.get('text', '')}\"", size=20, bold=False, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, Inches(2), Inches(3.2), Inches(6), Inches(0.4), f"— {sd.get('attribution', '')}", size=12, color=MUTED_TEXT, alignment=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
